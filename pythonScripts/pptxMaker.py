from pptx import Presentation
from pptx.util import Inches, Pt
import os

class tiredStudent():
    def __init__(self, slideData: dict, images: list) -> None:
        self.slideData = slideData
        self.images = images
        self.__prs = Presentation()
    
    def createPresentation(self) -> None:
        for i in range(1, len(self.slideData) + 1):
            self.__addSlide(str(i))
        self.__prs.save("public/slides/output.pptx")

    def __addSlide(self, key: str) -> None:
        slideObject = self.slideData[key]
        # title slide case, no body is given
        if len(slideObject["body"]) == 0:
            return self.__addTitleSlide(slideObject)
        # else just do a body picture slide,
        # TODO add more formats
        return self.__addBodyPictureSlide(key)

    def __addTitleSlide(self, slideObject: dict) -> None:
        slide = self.__prs.slides.add_slide(self.__prs.slide_layouts[0])
        slide.shapes.title.text = slideObject["title"]
        slide.placeholders[1].text = "Made With AI and Love"

    def __addBodyPictureSlide(self, key: str) -> None:
        slideObject = self.slideData[key]
        slide = self.__prs.slides.add_slide(self.__prs.slide_layouts[3])
        slide.shapes.title.text = slideObject["title"]
        tf = slide.shapes.placeholders[1].text_frame
        
        with open('image.png', 'wb') as handler:
            handler.write(self.images[int(key) - 1])
        slide.shapes.add_picture("image.png", Inches(5.6), Inches(2.5))
        if os.path.exists("image.png"):
            os.remove("image.png")
        
        for line in slideObject["body"]:
            p = tf.add_paragraph()
            p.font.size = Pt(16)
            if "• " in line:
                p.text = line[line.index("• ") + 2:]
                p.level = 1
            else:
                p.text = line
                p.level = 0