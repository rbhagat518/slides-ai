from pptx import Presentation
from pptx.util import Inches
import json

def addSlideObjectIntoSlideshow(key):
    slideObject = slides[key]
    # title slide case, no body is given
    if len(slideObject["body"]) == 0:
       slide = prs.slides.add_slide(prs.slide_layouts[0])
       slide.shapes.title.text = slideObject["title"]
       slide.placeholders[1].text = "Made With AI and Love"
       return

    # creating a slide that has a body
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = slideObject["title"]
    tf = slide.shapes.placeholders[1].text_frame
    slide.shapes.add_picture("public/slideComponents/images/image" + str(int(key) - 1) + ".png", Inches(5.6), Inches(2.5))
    for line in slideObject["body"]:
        p = tf.add_paragraph()
        if "â€¢ " in line:
            p.text = line[line.index("â€¢ ") + 4:]
            p.level = 1
        else:
            p.text = line
            p.level = 0

with open("public/slideComponents/slidesObjects/slides.json", "r" ) as f:
    slides = json.load(f)
prs = Presentation()
for i in range(1, len(slides) + 1):
    addSlideObjectIntoSlideshow(str(i))
prs.save("public/slideComponents/slides/output.pptx")