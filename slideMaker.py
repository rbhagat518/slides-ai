from pptx import Presentation
# script to parse the text file that the ai made
# and then convert it into a ppptx formmat

# read the text file into a string array
def textIntoLines(filename):
    lines = []
    with open(filename, 'r') as f:
        for line in f:
            if line == "\n" or len(line) == 0:
                continue
            lines.append(line[:line.index("\n")])
    return lines

# make the slides into tuples of format (slideNumber, title, bodyLines)
def linesIntoSlideObjects(lines):
    slides = []
    index = 0
    while index < len(lines):
        # skip over the slide number
        index += 1
        if index >= len(lines):
            break
        
        # get the title of the current slide
        title = lines[index]

        # move to the body
        index += 1
        if index >= len(lines):
            break

        # get the body to the current slide
        bodyLines = []
        while True:
            # move on once the whole body was read
            if lines[index] == f"Slide {len(slides) + 2}:":
                break
            
            # add this line to the body
            bodyLines.append(lines[index])

            # move to the next body line
            index += 1
            if index >= len(lines):
                break

        # add this slide to the list
        slides.append((len(slides) + 1, title, bodyLines))    
    return slides

def addSlideObjectIntoSlideshow(prensation, slideObject):
    # title slide case, no body is given
    if len(slideObject[2]) == 0:
       slide = prs.slides.add_slide(prs.slide_layouts[0])
       slide.shapes.title.text = slideObject[1]
       slide.placeholders[1].text = "Made With AI and Love"
       return

    # creating a slide that has a body
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slideObject[1]
    tf = slide.shapes.placeholders[1].text_frame
    for line in slideObject[2]:
        p = tf.add_paragraph()
        if "â€¢ " in line:
            p.text = line[line.index("â€¢ ") + 4:]
            p.level = 1
        else:
            p.text = line
            p.level = 0

prs = Presentation()
for slide in linesIntoSlideObjects(textIntoLines("inputTxt.txt")):
    addSlideObjectIntoSlideshow(prs, slide)
prs.save("output.pptx")